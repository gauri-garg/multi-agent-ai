from fastapi import FastAPI, Form, UploadFile, File, Depends, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from typing import Optional
from pydantic import BaseModel
from PIL import Image
import pytesseract
import io
import uuid
import json
import asyncio
import re
from sqlalchemy.orm import Session
from dotenv import load_dotenv

from agents.planner import plan_task
from agents.researcher import research_step
from agents.ml_agent import analyze_task
from memory.vector_db import store_memory, search_memory, load_db
from database import engine, Base, User, Chat, Message, SharedChat, SessionLocal
from auth import get_db, get_password_hash, verify_password, create_access_token, get_current_user, timedelta, ACCESS_TOKEN_EXPIRE_MINUTES

# 🔥 SAFE DL IMPORT
try:
    from agents.dl_agent import deep_analyze, load_dl_model
    DL_AVAILABLE = True
except:
    DL_AVAILABLE = False

import time
# Create DB Tables
max_retries = 3
for attempt in range(max_retries):
    try:
        Base.metadata.create_all(bind=engine)
        print("✅ DB connected")
        break
    except Exception as e:
        print(f"⚠️ DB connection failed (attempt {attempt + 1}/{max_retries}):", e)
        if attempt < max_retries - 1:
            time.sleep(3)  # Wait 3 seconds for Neon to wake up before retrying

app = FastAPI()


# ✅ CORS FIX
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class UserCreate(BaseModel):
    name: str
    email: str
    password: str

class UserLogin(BaseModel):
    email: str
    password: str

# ================= STARTUP =================
@app.on_event("startup")
def startup():
    load_db()
    if DL_AVAILABLE:
        load_dl_model()

@app.get("/")
def home():
    return {"message": "AI Running 🚀"}

# ================= AUTH ROUTES =================
@app.post("/auth/register")
def register(user: UserCreate, db: Session = Depends(get_db)):
    if db.query(User).filter(User.email == user.email).first():
        raise HTTPException(status_code=400, detail="Email already registered")
    new_user = User(name=user.name, email=user.email, password_hash=get_password_hash(user.password))
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    token = create_access_token(data={"sub": new_user.email}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    return {"access_token": token, "token_type": "bearer", "user": {"id": new_user.id, "name": new_user.name, "email": new_user.email}}

@app.post("/auth/login")
def login(user: UserLogin, db: Session = Depends(get_db)):
    db_user = db.query(User).filter(User.email == user.email).first()
    if not db_user or not verify_password(user.password, db_user.password_hash):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    token = create_access_token(data={"sub": db_user.email}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    return {"access_token": token, "token_type": "bearer", "user": {"id": db_user.id, "name": db_user.name, "email": db_user.email}}

@app.get("/auth/me")
def get_me(current_user: User = Depends(get_current_user)):
    return {"id": current_user.id, "name": current_user.name, "email": current_user.email}

@app.delete("/auth/delete-account")
def delete_account(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    db.delete(current_user)
    db.commit()
    return {"message": "Account deleted successfully"}

@app.get("/users")
def get_all_users(db: Session = Depends(get_db)):
    users = db.query(User).all()
    return [{"id": u.id, "name": u.name, "email": u.email, "created_at": u.created_at} for u in users]

# ================= SHARE ROUTES =================
@app.post("/chat/{chat_id}/share")
def share_chat(chat_id: int, is_public: bool = Form(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == current_user.id).first()
    if not chat: raise HTTPException(status_code=404)
    
    shared = db.query(SharedChat).filter(SharedChat.chat_id == chat.id).first()
    if not shared:
        shared = SharedChat(chat_id=chat.id, share_token=str(uuid.uuid4())[:12], is_public=is_public)
        db.add(shared)
    else:
        shared.is_public = is_public
    db.commit()
    return {"share_token": shared.share_token, "is_public": shared.is_public}

@app.get("/share/{token}")
def get_shared_chat(token: str, db: Session = Depends(get_db)):
    shared = db.query(SharedChat).filter(SharedChat.share_token == token, SharedChat.is_public == True).first()
    if not shared: raise HTTPException(status_code=404, detail="Share link invalid or private")
    
    messages = db.query(Message).filter(Message.chat_id == shared.chat_id).order_by(Message.id).all()
    msg_list = [{"id": m.id, "type": m.type, "text": m.text, "imageUrl": m.image_url, "response": m.response, "plan": json.loads(m.plan) if m.plan else [], "research": json.loads(m.research) if m.research else []} for m in messages]
    return {"name": shared.chat.name, "messages": msg_list}

# ================= CHAT =================
def serialize_chat(c):
    msgs = [{"id": m.id, "type": m.type, "text": m.text, "imageUrl": m.image_url, "response": m.response, "plan": json.loads(m.plan) if m.plan else [], "research": json.loads(m.research) if m.research else []} for m in c.messages]
    return {"id": c.id, "name": c.name, "pinned": c.pinned, "messages": msgs}

@app.get("/chats")
def get_chats(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    chats = db.query(Chat).filter(Chat.user_id == current_user.id).order_by(Chat.created_at.desc()).all()
    return [serialize_chat(c) for c in chats]

@app.post("/chat")
def create_chat(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    new_chat = Chat(user_id=current_user.id, name="New Chat")
    db.add(new_chat)
    db.commit()
    db.refresh(new_chat)
    return serialize_chat(new_chat)

@app.delete("/chat/{chat_id}")
def delete_chat(chat_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == current_user.id).first()
    if chat:
        db.delete(chat)
        db.commit()
    return {"message": "deleted"}

@app.put("/chat/{chat_id}/pin")
def pin_chat(chat_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == current_user.id).first()
    if chat:
        chat.pinned = not chat.pinned
        db.commit()
    return {"message": "toggled"}

# ================= RENAME CHAT =================
@app.put("/chat/{chat_id}/rename")
def rename_chat(chat_id: int, name: str = Form(...), current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == current_user.id).first()
    if chat:
        chat.name = name
        db.commit()
    return {"message": "renamed", "name": name}

# ================= DELETE MESSAGE (For Regeneration) =================
@app.delete("/chat/{chat_id}/message/{msg_id}")
def delete_message(chat_id: int, msg_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    msg = db.query(Message).join(Chat).filter(Message.id == msg_id, Chat.user_id == current_user.id).first()
    if msg:
        db.delete(msg)
        db.commit()
    return {"message": "deleted"}

# ================= TEXT AND FILE UPLOAD =================
@app.post("/chat/{chat_id}")
async def add_message(
    chat_id: int, 
    task: str = Form(""), 
    file: Optional[UploadFile] = File(None),
    ai_length: str = Form("Auto"),
    ai_format: str = Form("Auto"),
    ai_tone: str = Form("Auto"),
    ai_language: str = Form("Auto"),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    chat = db.query(Chat).filter(Chat.id == chat_id, Chat.user_id == current_user.id).first()
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")

    extracted_text = ""
    img_data_url = None

    if file and file.filename:
        filename = file.filename.lower()
        contents = await file.read()
        if not contents:
            raise HTTPException(status_code=400, detail="Empty file uploaded")

        if filename.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif")):
            image = Image.open(io.BytesIO(contents))
            extracted_text = pytesseract.image_to_string(image)
            try:
                preview_img = image.copy()
                preview_img.thumbnail((500, 500))
                if preview_img.mode in ("RGBA", "P"):
                    preview_img = preview_img.convert("RGB")
                buffered = io.BytesIO()
                preview_img.save(buffered, format="JPEG")
                import base64
                base64_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
                img_data_url = f"data:image/jpeg;base64,{base64_str}"
            except Exception:
                img_data_url = None
        elif filename.endswith(".pdf"):
            try:
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text: extracted_text += text + "\n"
            except ImportError:
                raise HTTPException(status_code=500, detail="PyPDF2 is not installed. Run `pip install PyPDF2`")
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"PDF Extraction Error: {e}")
        elif filename.endswith(".docx"):
            try:
                import docx
                doc = docx.Document(io.BytesIO(contents))
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                raise HTTPException(status_code=500, detail="python-docx is not installed. Run `pip install python-docx`")
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"DOCX Extraction Error: {e}")
        elif filename.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(contents))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            except ImportError:
                raise HTTPException(status_code=500, detail="python-pptx is not installed. Run `pip install python-pptx`")
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"PPTX Extraction Error: {e}")
        elif filename.endswith((".txt", ".csv", ".json", ".md")):
            extracted_text = contents.decode("utf-8", errors="ignore")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")

        if not extracted_text.strip():
            extracted_text = f"[Uploaded file: {file.filename}, but no text was found.]"
        
        # Groq free tier limit is 6000 tokens/min. 8000 chars is ~2000 tokens.
        if len(extracted_text) > 8000:
            extracted_text = extracted_text[:8000] + "\n\n[Content truncated due to size limits...]"

    final_prompt = task.strip()
    if file and file.filename:
        if final_prompt:
            final_prompt = f"User Prompt: {final_prompt}\n\n--- ATTACHED FILE CONTENT ({file.filename}) ---\n{extracted_text}\n--- END OF FILE CONTENT ---"
        else:
            final_prompt = f"Analyze and summarize this attached file ({file.filename}):\n\n--- FILE CONTENT ---\n{extracted_text}\n--- END OF FILE CONTENT ---"

    if not final_prompt:
        raise HTTPException(status_code=400, detail="No prompt or file provided")

    display_text = task.strip()
    if file and file.filename:
        file_badge = f"[File Uploaded: {file.filename} 📎]"
        display_text = f"{file_badge}\n{display_text}".strip()

    user_msg = Message(chat_id=chat.id, type="user", text=display_text, image_url=img_data_url)
    db.add(user_msg)
    db.commit()

    # Build Context
    messages_dict = [{"type": m.type, "text": m.text, "response": m.response, "plan": json.loads(m.plan) if m.plan else []} for m in chat.messages]
    context = build_smart_context(messages_dict)

    async def stream_response():
        ai_msg = await asyncio.to_thread(generate_ai_response, final_prompt, context, messages_dict, ai_length, ai_format, ai_tone, ai_language)
        text = str(ai_msg.get("response", ""))
        tokens = re.split(r'(\s+)', text)
        for token in tokens:
            if token:
                yield json.dumps({"type": "chunk", "text": token}) + "\n"
                await asyncio.sleep(0.02)  # Fast, realistic stream delay
        
        # Save AI Response
        ai_msg_db = Message(
            chat_id=chat.id, type="ai", response=ai_msg.get("response"),
            plan=json.dumps(ai_msg.get("plan", [])), research=json.dumps(ai_msg.get("full_research", []))
        )
        db_session = SessionLocal()
        try:
            db_session.add(ai_msg_db)
            db_session.commit()
            db_session.refresh(ai_msg_db)
            ai_msg["id"] = ai_msg_db.id
        finally:
            db_session.close()
        
        yield json.dumps({"type": "done", "full": ai_msg}) + "\n"

    return StreamingResponse(stream_response(), media_type="application/x-ndjson")

##build_context(task)
def build_smart_context(messages, limit=6):
    context = []

    for msg in messages[-limit:]:
        if msg["type"] == "user":
            # Truncate user messages in context to save tokens
            short_user = msg['text'][:300] + "..." if len(msg['text']) > 300 else msg['text']
            context.append(f"User: {short_user}")
        elif msg["type"] == "ai":
            short = msg.get("response", "")[:120]
            context.append(f"AI: {short}")

    return "\n".join(context)

def detect_query_type(task, context):
    task_lower = task.lower()

    if task_lower in context.lower():
        return "repeat"

    if any(word in task_lower for word in ["summary", "short", "brief"]):
        return "summary"

    if any(word in task_lower for word in ["more", "explain", "detail"]):
        return "followup"

    return "new"

# ================= SMART FORMAT CLASSIFIER =================
def detect_format_type(query):
    query_lower = query.lower()
    if "compare" in query_lower or "difference" in query_lower:
        return "Table"
    elif "steps" in query_lower or "how to" in query_lower or "guide" in query_lower:
        return "Step-by-step"
    elif "list" in query_lower or "features" in query_lower:
        return "Bullet Points"
    elif "flowchart" in query_lower or "process" in query_lower or "decision" in query_lower:
        return "Flowchart"
    else:
        return "Paragraph"

# ================= AI CORE =================
def generate_ai_response(task: str, context="", messages=[], ai_length="Auto", ai_format="Auto", ai_tone="Auto", ai_language="Auto"):

    # 🔥 SMART TYPE DETECTION
    query_type = detect_query_type(task, context)
    
    # Truncate for classifiers to avoid memory blowout with large files
    short_task = task[:1000]

    # =============================
    # 1. PLAN
    # =============================
    planning_prompt = f"Break down this task into 3-5 concise search queries or sub-topics for research: {short_task}"
    steps = plan_task(planning_prompt)

    if isinstance(steps, list):
        steps_list = steps
    else:
        steps_list = [s.strip() for s in steps.split("\n") if s.strip()]

    steps_list = steps_list[:6]

    # =============================
    # 2. ML + DL
    # =============================
    ml_result = analyze_task(short_task)

    if DL_AVAILABLE:
        dl_result = deep_analyze(short_task)
    else:
        dl_result = {"dl_category": "N/A", "dl_confidence": "0%"}

    # =============================
    # 3. MEMORY
    # =============================
    store_memory(short_task, metadata={"type": "task"})
    memory = search_memory(short_task)

    # =============================
    # 4. RESEARCH (SMART FILTER)
    # =============================
    existing_steps = []
    for m in messages:
        if m.get("plan"):
            existing_steps.extend(m["plan"])

    research_results = []

    for step in steps_list:
        if step in existing_steps:
            continue  # 🔥 skip repeated work

        research_results.append({
            "step": step,
            "research": research_step(step)
        })

    research_results = research_results[:3]
    research_text = "\n".join([f"- {r['step']}: {r['research'][:1000]}" for r in research_results])

    # =============================
    # 5. STRICT PROMPT BUILDER & AI SYNTHESIS
    # =============================
    detected_format = None
    if ai_format == "Auto":
        detected_format = detect_format_type(task)
        format_instruction = f"""
Analyze the user query and choose the BEST format:
- Use bullet points → for lists or features
- Use step-by-step → for guides/tutorials
- Use table → for comparisons or structured data
- Use flowchart → for processes or decision making
- Otherwise → use paragraph

Smart Detection Suggestion: {detected_format}
"""
    else:
        format_instruction = f"Respond strictly in {ai_format} format."

    system_instruction = f"""
You are an elite, highly capable AI assistant (like ChatGPT / Gemini). Your intelligence is vast.

User Preferences:
- Tone: {ai_tone}
- Mode: {ai_length}
- Format: {ai_format}
- Language: {ai_language}

Rules:
{format_instruction}

Tone Rules:
- Professional → formal and precise
- Casual → friendly and simple
- Technical → detailed and domain-specific
- Friendly → warm and easy to read
- Auto → adopt a natural, conversational tone

Mode Rules:
- Summary → short and concise
- Detailed → deep explanation
- Auto → decide based on query complexity

Formatting Rules:
- Bullet Points → EACH point MUST be on a NEW LINE starting with a • symbol. Do NOT add empty lines between bullet points.
- Step-by-step → EACH step MUST be on a NEW LINE with a clear number (1., 2., etc). Do NOT add empty lines between steps.
- Table → Return a proper, complete markdown table. Do NOT break the columns.
- Flowchart → Generate a valid Mermaid.js code block (```mermaid ... ```). Nodes MUST use simple alphanumeric IDs without spaces.
- Paragraph → readable paragraphs with spacing

CRITICAL:
- If user selected a format (not Auto) → DO NOT change it.
- If format is Auto → choose the best format intelligently.
- Response must be clean, readable, and structured. Avoid excessive line spacing or blank lines.
- Answer directly using the Knowledge Base. Remove robotic filler text.
   - NEVER generate repeated headings or duplicated paragraphs.
   - NEVER repeat the exact same sentence twice.
DO NOT:
- Do NOT merge separate bullet points onto a single line.
- Do NOT write everything in one giant paragraph.
- Do NOT ignore the requested User Format.

User Format: {ai_format}


Knowledge Base:
{research_text}

User Query:
{task}
"""

    # Send the combined data back to the LLM to generate the TRUE final response
    final_response = plan_task(system_instruction)
    
    if isinstance(final_response, list):
        final_response = "\n".join(str(x) for x in final_response)

    # =============================
    # 6. RETRY ENFORCEMENT LOGIC
    # =============================
    if ai_format == "Table" and "|" not in final_response:
        print("⚠️ AI ignored Table format. Retrying...")
        retry_prompt = system_instruction + "\n\n[SYSTEM ERROR]: You failed validation. Rewrite your previous response entirely as a COMPLETE MARKDOWN TABLE."
        final_response = plan_task(retry_prompt)
        if isinstance(final_response, list):
            final_response = "\n".join(str(x) for x in final_response)
    elif ai_format == "Bullet Points" and not any(c in final_response for c in ["- ", "* ", "• "]):
        print("⚠️ AI ignored Bullet format. Retrying...")
        retry_prompt = system_instruction + "\n\n[SYSTEM ERROR]: You failed validation. Rewrite your previous response entirely as a BULLETED LIST."
        final_response = plan_task(retry_prompt)
        if isinstance(final_response, list):
            final_response = "\n".join(str(x) for x in final_response)

    # Append Metrics
    final_response += f"\n\n🧠 ML: {ml_result['category']} ({ml_result['confidence']})"
    final_response += f"\n🤖 DL: {dl_result['dl_category']} ({dl_result['dl_confidence']})"

    return {
        "type": "ai",
        "response": final_response,
        "plan": steps_list,
        "research": research_results[:2],
        "full_research": research_results,
        "memory": memory,
        "detected_format": detected_format if ai_format == "Auto" else None
    }